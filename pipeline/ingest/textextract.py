"""A3: text extraction per format.

Every extractor returns an ExtractResult:
    method:  extract_method label (contract §4.1)
    pages:   list[(page_no:int, text:str)]  -- 1-based logical pages
    n_pages: int
    ocr_pages: list[int]  -- page numbers whose text layer is too thin (pdf only)
"""
from __future__ import annotations

import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from . import config
from .util import log


@dataclass
class ExtractResult:
    method: str
    pages: list[tuple[int, str]]
    n_pages: int
    ocr_pages: list[int] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


# --------------------------------------------------------------------------- PDF
def extract_pdf(path: Path) -> ExtractResult:
    import fitz  # PyMuPDF

    pages: list[tuple[int, str]] = []
    ocr_pages: list[int] = []
    warnings: list[str] = []
    with fitz.open(path) as doc:
        n = doc.page_count
        for i in range(n):
            try:
                txt = doc.load_page(i).get_text("text")
            except Exception as e:  # noqa: BLE001
                txt = ""
                warnings.append(f"page {i+1} read error: {e}")
            pages.append((i + 1, txt))
            if len(txt.strip()) < config.OCR_MIN_TEXT_CHARS:
                ocr_pages.append(i + 1)
    return ExtractResult("pymupdf", pages, len(pages), ocr_pages, warnings)


# -------------------------------------------------------------------------- DOCX
def extract_docx(path: Path) -> ExtractResult:
    import docx  # python-docx

    try:
        d = docx.Document(str(path))
    except Exception:
        # .docm (macro-enabled) may carry a content type python-docx rejects;
        # fall back to parsing word/document.xml directly from the zip.
        return _extract_docx_xml(path)
    parts: list[str] = []
    for para in d.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    text = "\n".join(parts)
    return ExtractResult("docx", [(1, text)], 1)


def _extract_docx_xml(path: Path) -> ExtractResult:
    """Raw OOXML fallback: read word/document.xml, join runs per paragraph."""
    import zipfile
    from lxml import etree

    W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    with zipfile.ZipFile(path) as zf:
        data = zf.read("word/document.xml")
    root = etree.fromstring(data)
    paras: list[str] = []
    for p in root.iter(W + "p"):
        texts = [t.text or "" for t in p.iter(W + "t")]
        s = "".join(texts).strip()
        if s:
            paras.append(s)
    text = "\n".join(paras)
    return ExtractResult("docx", [(1, text)], 1,
                         warnings=["docm parsed via raw OOXML fallback"])


# --------------------------------------------------------------------------- DOC
def extract_doc(path: Path) -> ExtractResult:
    """Legacy .doc via macOS textutil, fallback to soffice->txt."""
    warnings: list[str] = []
    # 1) textutil (macOS)
    try:
        proc = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, timeout=120,
        )
        if proc.returncode == 0 and proc.stdout.strip():
            text = proc.stdout.decode("utf-8", "replace")
            return ExtractResult("doc", [(1, text)], 1, warnings=warnings)
        warnings.append("textutil produced no text")
    except (OSError, subprocess.TimeoutExpired) as e:
        warnings.append(f"textutil error: {e}")

    # 2) soffice fallback
    text = _soffice_to_txt(path, warnings)
    if text is not None:
        return ExtractResult("doc", [(1, text)], 1, warnings=warnings)
    raise RuntimeError("doc extraction failed: " + "; ".join(warnings))


def _soffice_to_txt(path: Path, warnings: list[str]) -> str | None:
    import shutil
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        warnings.append("soffice not available")
        return None
    with tempfile.TemporaryDirectory() as tmp:
        try:
            proc = subprocess.run(
                [soffice, "--headless", "--convert-to", "txt:Text",
                 "--outdir", tmp, str(path)],
                capture_output=True, timeout=180,
            )
        except (OSError, subprocess.TimeoutExpired) as e:
            warnings.append(f"soffice error: {e}")
            return None
        if proc.returncode != 0:
            warnings.append(f"soffice rc={proc.returncode}")
            return None
        out = Path(tmp) / (path.stem + ".txt")
        if out.exists():
            return out.read_text("utf-8", "replace")
        warnings.append("soffice produced no txt")
        return None


# -------------------------------------------------------------------------- PPTX
def extract_pptx(path: Path) -> ExtractResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[tuple[int, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
        pages.append((idx, "\n".join(parts)))
    return ExtractResult("pptx", pages, len(pages))


# ---------------------------------------------------------------------- XLS/XLSX
def extract_xlsx(path: Path) -> ExtractResult:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages: list[tuple[int, str]] = []
    for idx, ws in enumerate(wb.worksheets, start=1):
        lines = [f"# sheet: {ws.title}"]
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                lines.append("\t".join(vals).rstrip())
        pages.append((idx, "\n".join(lines)))
    wb.close()
    return ExtractResult("xlsx", pages, len(pages))


def extract_xls(path: Path) -> ExtractResult:
    import xlrd

    book = xlrd.open_workbook(str(path))
    pages: list[tuple[int, str]] = []
    for idx in range(book.nsheets):
        sh = book.sheet_by_index(idx)
        lines = [f"# sheet: {sh.name}"]
        for r in range(sh.nrows):
            vals = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
            if any(v.strip() for v in vals):
                lines.append("\t".join(vals).rstrip())
        pages.append((idx + 1, "\n".join(lines)))
    return ExtractResult("xls", pages, len(pages))


DISPATCH = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".docm": extract_docx,   # docm is docx-structured
    ".doc": extract_doc,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls": extract_xls,
}


def extract(path: Path) -> ExtractResult:
    ext = path.suffix.lower()
    fn = DISPATCH.get(ext)
    if fn is None:
        raise ValueError(f"no extractor for {ext}")
    return fn(path)
